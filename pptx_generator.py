import io
from datetime import datetime

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# WMO weather code → (emoji, Japanese name)
WMO_MAP: dict[int, tuple[str, str]] = {
    0:  ("☀",  "快晴"),
    1:  ("🌤", "晴れ"),
    2:  ("⛅", "晴れ時々曇り"),
    3:  ("☁",  "曇り"),
    45: ("🌫", "霧"),
    48: ("🌫", "霧(着氷)"),
    51: ("🌦", "霧雨(弱)"),
    53: ("🌦", "霧雨"),
    55: ("🌦", "霧雨(強)"),
    56: ("🌧", "凍雨(弱)"),
    57: ("🌧", "凍雨"),
    61: ("🌧", "雨(弱)"),
    63: ("🌧", "雨"),
    65: ("🌧", "雨(強)"),
    66: ("🌧", "凍雨(弱)"),
    67: ("🌧", "凍雨(強)"),
    71: ("❄",  "雪(弱)"),
    73: ("❄",  "雪"),
    75: ("❄",  "雪(強)"),
    77: ("❄",  "霧雪"),
    80: ("🌦", "にわか雨(弱)"),
    81: ("🌦", "にわか雨"),
    82: ("🌦", "にわか雨(強)"),
    85: ("🌨", "にわか雪(弱)"),
    86: ("🌨", "にわか雪(強)"),
    95: ("⛈",  "雷雨"),
    96: ("⛈",  "雷雨(雹)"),
    99: ("⛈",  "激しい雷雨"),
}

WEEKDAY_JA = ["月", "火", "水", "木", "金", "土", "日"]

# Colors
COLOR_HEADER_DATE = RGBColor(0xBD, 0xD7, 0xEE)   # light blue
COLOR_HEADER_LOC  = RGBColor(0xD9, 0xD9, 0xD9)   # light gray
COLOR_LOC_CELL    = RGBColor(0xFF, 0xE5, 0xCC)   # light orange
COLOR_BORDER      = RGBColor(0x80, 0x80, 0x80)   # gray
COLOR_CREDIT      = RGBColor(0x80, 0x80, 0x80)
COLOR_SCHEDULE_BG = RGBColor(0xFF, 0xFF, 0xFF)


def _wmo(code) -> tuple[str, str]:
    if code is None:
        return ("❓", "不明")
    return WMO_MAP.get(int(code), ("❓", "不明"))


def _fmt_date(date_str: str) -> str:
    d = datetime.strptime(date_str, "%Y-%m-%d")
    wd = WEEKDAY_JA[d.weekday()]
    return f"{d.month}/{d.day}({wd})"


def _set_cell_bg(cell, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")


def _add_paragraph(tf, text: str, font_size: int, bold: bool = False,
                   align=PP_ALIGN.CENTER, color: RGBColor | None = None):
    p = tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def generate_pptx(
    locations: list[dict],
    forecasts: list[dict],
    days: int,
) -> bytes:
    """
    locations: [{name, lat, lon}, ...]
    forecasts: [daily_dict, ...]  (same order as locations)
    days: requested forecast days
    Returns .pptx as bytes.
    """
    prs = Presentation()
    prs.slide_width  = Cm(33.867)
    prs.slide_height = Cm(19.05)

    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # ── Schedule area (top 1/3) ──────────────────────────────────────
    schedule_box = slide.shapes.add_textbox(
        Cm(0.5), Cm(0.3), Cm(32.867), Cm(5.75)
    )
    schedule_box.text_frame.word_wrap = True
    tf = schedule_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "予定欄"
    run.font.size = Pt(9)
    run.font.bold = False
    run.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

    # Border for schedule box
    from pptx.oxml.ns import qn
    from lxml import etree

    sp = schedule_box._element
    spPr = sp.find(qn("p:spPr"))
    ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", "12700")  # 1pt border
    solidFill = etree.SubElement(ln, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", "AAAAAA")

    # ── Weather table ────────────────────────────────────────────────
    # Determine actual days available (may be less than requested)
    actual_days = days
    if forecasts:
        actual_days = min(days, len(forecasts[0].get("time", [])))

    n_rows = 1 + len(locations)          # header + location rows
    n_cols = 1 + actual_days             # location name col + date cols

    table_left   = Cm(0.5)
    table_top    = Cm(6.4)
    table_width  = Cm(32.867)
    table_height = Cm(12.0)

    tbl = slide.shapes.add_table(n_rows, n_cols, table_left, table_top,
                                 table_width, table_height).table

    # Column widths
    loc_col_width  = Cm(3.0)
    date_col_width = (table_width - loc_col_width) // actual_days

    tbl.columns[0].width = loc_col_width
    for c in range(1, n_cols):
        tbl.columns[c].width = date_col_width

    # Row heights
    header_row_height = Cm(0.8)
    data_row_height   = (table_height - header_row_height) // len(locations)

    tbl.rows[0].height = header_row_height
    for r in range(1, n_rows):
        tbl.rows[r].height = data_row_height

    # ── Header row ───────────────────────────────────────────────────
    cell = tbl.cell(0, 0)
    _set_cell_bg(cell, COLOR_HEADER_LOC)
    tf = cell.text_frame
    tf.clear()
    _add_paragraph(tf, "地点", 9, bold=True)

    times = forecasts[0].get("time", []) if forecasts else []
    for c in range(1, n_cols):
        cell = tbl.cell(0, c)
        _set_cell_bg(cell, COLOR_HEADER_DATE)
        tf = cell.text_frame
        tf.clear()
        date_label = _fmt_date(times[c - 1]) if (c - 1) < len(times) else ""
        _add_paragraph(tf, date_label, 8, bold=True)

    # ── Data rows ─────────────────────────────────────────────────────
    for r, (loc, daily) in enumerate(zip(locations, forecasts), start=1):
        loc_name = loc["name"][:10]  # truncate long names

        # Location name cell
        cell = tbl.cell(r, 0)
        _set_cell_bg(cell, COLOR_LOC_CELL)
        tf = cell.text_frame
        tf.clear()
        tf.word_wrap = True
        _add_paragraph(tf, loc_name, 8, bold=True)

        # Weather cells
        codes  = daily.get("weather_code", [])
        maxts  = daily.get("temperature_2m_max", [])
        mints  = daily.get("temperature_2m_min", [])
        precip = daily.get("precipitation_probability_max", [])

        for c in range(1, n_cols):
            i = c - 1
            cell = tbl.cell(r, c)
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = False

            emoji, name = _wmo(codes[i] if i < len(codes) else None)
            maxt  = f"{maxts[i]:.0f}" if i < len(maxts) and maxts[i] is not None else "--"
            mint  = f"{mints[i]:.0f}" if i < len(mints) and mints[i] is not None else "--"
            prec  = f"{int(precip[i])}%" if (i < len(precip) and precip[i] is not None) else "-"

            # Line 1: weather
            _add_paragraph(tf, f"{emoji} {name}", 7)
            # Line 2: precipitation
            _add_paragraph(tf, f"降水 {prec}", 7)
            # Line 3: temp
            _add_paragraph(tf, f"{maxt}°/ {mint}°", 7)

    # ── Credit text ───────────────────────────────────────────────────
    credit = slide.shapes.add_textbox(
        Cm(22.0), Cm(18.4), Cm(11.5), Cm(0.5)
    )
    tf = credit.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "出典: Open-Meteo (CC BY 4.0)"
    run.font.size = Pt(7)
    run.font.color.rgb = COLOR_CREDIT

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
